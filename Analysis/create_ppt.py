
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "F1 Insights & Strategy Explorer"
    subtitle.text = "Group OOAD 23\nMembers: Abhi Bhardwaj & Ankita Pimpalkar"

    # Slide 2: Project Description
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Project Description"
    content = slide.placeholders[1]
    content.text = ("The F1 Insights & Strategy Explorer is a comprehensive tool "
                    "designed for Formula 1 fans and analysts to visualize and interpret "
                    "complex race data. By selecting specific race sessions and drivers, "
                    "users can analyze performance through lap time comparisons, position "
                    "tracking throughout the race, and detailed tyre strategy (stint) "
                    "visualizations. The system integrates data from external sources "
                    "like the FastF1 API to provide accurate and up-to-date insights "
                    "into race strategy.")

    # Slide 3: Use Case Model
    slide_layout = prs.slide_layouts[5] # Title Only (Blank frame)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Use Case Model"
    
    use_case_img = "/Users/abhiabhardwaj/Analysis/UseCaseDiagram.png"
    if os.path.exists(use_case_img):
        # Add image centered
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(use_case_img, left, top, height=height)

    # Slides 4-8: Robustness Diagrams
    base_path = "/Users/abhiabhardwaj/Analysis/"
    slide_data = [
        ("Select Season/Race/Session", "Select_Season_Race_Session_Handwritten.png"),
        ("Select Drivers", "Select_Drivers_Handwritten.png"),
        ("View Lap Time Comparison", "View_Lap_Time_Comparison_Handwritten.png"),
        ("View Position Over Laps", "View_Position_Over_Laps_Handwritten.png"),
        ("View Tyre Stints", "View_Tyre_Stints_Handwritten.png"),
    ]

    for title_text, filename in slide_data:
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Add image
        img_path = base_path + filename
        left = Inches(1)
        top = Inches(2)
        height = Inches(4.5)
        slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save('/Users/abhiabhardwaj/Analysis/F1_Robustness_Analysis.pptx')
    print("Presentation saved to /Users/abhiabhardwaj/Analysis/F1_Robustness_Analysis.pptx")

create_presentation()
